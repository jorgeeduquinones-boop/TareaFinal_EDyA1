"""
Actualiza la presentación PPTX existente sin reescribirla:
  - Portada: equipo completo, grupo, códigos.
  - Slide nueva (al final): aplicativo web + URL del demo.
  - Slide nueva: soporte al formato extendido del profesor.
  - Sustenta '10 tests' → '11 tests'.

Preserva estilos del .pptx original (no crea uno nuevo).
"""
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

IN = "TareaFinal_EDyA1_Presentacion.pptx"
OUT = "TareaFinal_EDyA1_Presentacion.pptx"  # in-place

prs = Presentation(IN)

# -------------------------------------------------------------------- portada
slide1 = prs.slides[0]
# Shape 8 = línea con el autor; reemplazamos por el equipo completo.
# Shape 9 = entrega. Insertamos otra línea con grupo + códigos.
def set_text_preserve_style(shape, new_text):
    """Reemplaza el texto del shape conservando el formato del primer run."""
    tf = shape.text_frame
    p = tf.paragraphs[0]
    # Guarda formato del primer run
    if p.runs:
        r0 = p.runs[0]
        font_name = r0.font.name
        font_size = r0.font.size
        font_bold = r0.font.bold
        color_rgb = None
        try:
            color_rgb = r0.font.color.rgb
        except (AttributeError, Exception):
            color_rgb = None
    else:
        font_name, font_size, font_bold, color_rgb = None, None, None, None
    # Limpia párrafos extra
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    # Limpia runs
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    # Nuevo run con formato preservado
    nr = p.add_run()
    nr.text = new_text
    if font_name: nr.font.name = font_name
    if font_size: nr.font.size = font_size
    if font_bold is not None: nr.font.bold = font_bold
    if color_rgb: nr.font.color.rgb = color_rgb

# Equipo en lugar del solo autor (shape 8)
set_text_preserve_style(
    slide1.shapes[8],
    "Joan Mateo Cardona 2243431  •  Danna Villegas 2240027  •  Jorge Eduardo Álvarez 2230610"
)

# Línea de entrega ahora incluye el grupo (shape 9)
set_text_preserve_style(
    slide1.shapes[9],
    "Grupo 1  •  Entrega: 17 de mayo de 2026  •  Sustentación: semana 16"
)

# -------------------------------------------------- nueva slide: aplicativo web
def add_text_slide(prs, layout_idx_fallback=5):
    """Crea una slide nueva clonando el layout 'BASE' del slide 1 si es posible."""
    # Usar el mismo layout que el slide 1 para que el fondo coincida
    layout = prs.slides[0].slide_layout
    return prs.slides.add_slide(layout)

def add_titulo(slide, texto, *, top_in=0.6, left_in=0.8, width_in=11.5, height_in=0.8, size_pt=36, bold=True):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = texto
    r.font.size = Pt(size_pt)
    r.font.bold = bold
    r.font.name = "Calibri"
    r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    return tb

def add_sub(slide, texto, *, top_in, left_in=0.8, width_in=11.5, height_in=0.4, size_pt=18, color=(0x47, 0x55, 0x69)):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = texto
    r.font.size = Pt(size_pt)
    r.font.name = "Calibri"
    r.font.color.rgb = RGBColor(*color)
    return tb

def add_bullets(slide, items, *, top_in, left_in=0.8, width_in=11.5, height_in=5.0, size_pt=18):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = "• " + item
        r.font.size = Pt(size_pt)
        r.font.name = "Calibri"
        r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        p.space_after = Pt(8)
    return tb

# Slide nueva 1: aplicativo web entregado
s_app = add_text_slide(prs)
add_titulo(s_app, "Aplicativo web entregado")
add_sub(s_app, "PWA con cinco paneles + integración continua", top_in=1.4)
add_bullets(s_app, [
    "Resultados: ranking, métricas y gráficos del ejecutado.",
    "Cargar caso: pegado, CSV, generador sintético (m, n, p, seed) y casos del profesor.",
    "Benchmark: MergeSort, QuickSort 3-way, RadixSort, Insertion y Array.sort sobre el mismo dataset.",
    "Pruebas: suite automatizada de 11 tests (Node + navegador).",
    "Documentación: complejidad, decisiones de diseño y trazabilidad de prácticas.",
    "CI con GitHub Actions ejecuta los tests antes de cada despliegue.",
], top_in=2.1, size_pt=18)
add_sub(
    s_app,
    "Demo en vivo: https://tarea-final-e-dy-a1.vercel.app/app/index.html",
    top_in=6.7, size_pt=16, color=(0x25, 0x63, 0xEB),
)

# Slide nueva 2: formato extendido del profesor
s_fmt = add_text_slide(prs)
add_titulo(s_fmt, "Soporte al formato del profesor")
add_sub(s_fmt, "Casos publicados en data_2026_1/", top_in=1.4)
add_bullets(s_fmt, [
    "Formato extendido: 'Cat1,Cat2,…,CatN--customer product category price quantity timestamp;…'",
    "El parser detecta automáticamente comillas externas, separador '--' y lista de categorías.",
    "calcularPedidos sigue con la firma original: la normalización ocurre antes del scan.",
    "Caso 0 — 20 clientes · 7 categorías · 182 records.",
    "Caso 1 — 50 clientes · 20 categorías · 1 000 records.",
    "Caso 2 — 100 clientes · 30 categorías · 5 000 records.",
    "Tiempos medidos: 1 a 5 ms en ambas implementaciones (QuickSort y MergeSort).",
], top_in=2.1, size_pt=18)
add_sub(
    s_fmt,
    "Test 11 verifica que un caso con y sin prefijo de categorías produce el mismo ranking.",
    top_in=6.7, size_pt=16, color=(0x47, 0x55, 0x69),
)

prs.save(OUT)
print(f"OK: {OUT}")
